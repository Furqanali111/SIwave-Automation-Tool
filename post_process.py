import pandas as pd
from pptx.enum.text import PP_ALIGN
import matplotlib.colors as mcolors
import pptx
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
import os
from pandas.plotting import table
import dataframe_image as dfi
import numpy as np
import math
import openpyxl
from PIL import Image
import math
from milkviz import dot_heatmap
import seaborn as sns
from PIL import Image, ExifTags
import sys



# ## iput parameters ##

path0=sys.argv[1]
print("excel file path path ",path0)
df_parameters_input = pd.read_excel(path0)
df_parameters_input.head()

siwfolderpath=sys.argv[2]
print("siwfolderpath: ",siwfolderpath)

# ## input IMAX files ##

path = os.getcwdb()
path = path.decode("utf-8")
#print(path)
directory = path

# Extract the column names from the DataFrame
column_names = df_parameters_input.columns.tolist()
# List files ending with 'DCR.dc' in a specific folder/directory and also present as a column name in df_parameters_input
list_of_files_IMAX = [f for f in os.listdir(directory) if f.endswith('IMAX.dc') and os.path.splitext(f)[0].rstrip('_IMAX') in column_names]
# Add the directory path to the filenames
list_of_files_IMAX = [os.path.join(directory, f) for f in list_of_files_IMAX]
#print(list_of_files)
print(list_of_files_IMAX)



Imax_df = pd.DataFrame(columns=['Power Rail', 'Current (A)','Imax (A) Bump','Imax (A) BGA'])


# ## find bump map layout ##


def find_layout(current_file):
    target_line = "Vias"
    remove_keywords = "Voltage Probes"
    # boolean flag to track if the target line has been reached
    found = False
    # list to store extracted information
    data = []
    with open(current_file, "r") as f:
        for line in f:
            # set found to True if the target line is reached
            if target_line in line:
                break
            # add line to data list if the target line has been reached
        for line in f:
            include_line = True
            if remove_keywords in line:
                break
            if include_line and line.strip():
                data.append(line.strip().split('\t'))

    columns = data[0]
    data = data[1:]
    df_current = pd.DataFrame(data, columns=columns)
    return df_current


# In[21]:


layout = pd.DataFrame(columns = ['Via', 'Net', 'x (mm)', 'y (mm)', 'Current / A', 'Limit / A', 'Pass / Fail', 'Resistance / Ohms', 'IR Drop / V', 'Power / W'])


# In[22]:


for file_name in list_of_files_IMAX:
    current_file = file_name
    df_current = find_layout(current_file)
    layout = pd.concat([layout, df_current], axis=0)

bump_layout = layout[layout["Via"].str.contains('C4 Bump')]
bump_layout['x_cp (mm)'] = pd.to_numeric(bump_layout['x (mm)'])
bump_layout['y_cp (mm)'] = pd.to_numeric(bump_layout['y (mm)'])
bump_max_x = math.floor(bump_layout['x_cp (mm)'].max() * 2) / 2
bump_min_x = math.floor(bump_layout['x_cp (mm)'].min() * 2) / 2
bump_max_y = math.floor(bump_layout['y_cp (mm)'].max() * 2) / 2
bump_min_y = math.floor(bump_layout['y_cp (mm)'].min() * 2) / 2

bump_layout['Current_copy'] = pd.to_numeric(bump_layout['Current / A']).abs()
# bump_min_current = math.floor(bump_layout['Current_copy'].min() * 2) / 2
# bump_max_current = math.floor(bump_layout['Current_copy'].max() * 2) / 2
# # bump_layout = bump_layout.sort_values(by=['x_cp (mm)'], ascending=False)
# print('min_current:', bump_min_current)

# Define the custom function for multiplication
def multiply_with_params(row, params_df):
    net = row['Net']
    current_copy = row['Current_copy']
    if net in params_df.columns:
        return current_copy * params_df.loc[0, net]
    return current_copy

# Multiply each row and store the results in a new column 'EDC_adjusted'
bump_layout['EDC_adjusted'] = bump_layout.apply(lambda row: multiply_with_params(row, df_parameters_input), axis=1)

# Find the minimum and maximum values in the 'EDC_adjusted' column
min_edc_adjusted = bump_layout['EDC_adjusted'].min()
max_edc_adjusted = bump_layout['EDC_adjusted'].max()

# Convert to float variables
min_edc_adjusted = float(min_edc_adjusted)
max_edc_adjusted = float(max_edc_adjusted)

print("Minimum value in 'EDC_adjusted':", min_edc_adjusted)
print("Maximum value in 'EDC_adjusted':", max_edc_adjusted)

bump_xratio = int(bump_max_x - bump_min_x)
bump_yratio = int(bump_max_y - bump_min_y)

# ## find BGA map layout ##
BGA_layout = layout[layout["Via"].str.contains('Ball')]
BGA_layout['x_cp (mm)'] = pd.to_numeric(BGA_layout['x (mm)'])
BGA_layout['y_cp (mm)'] = pd.to_numeric(BGA_layout['y (mm)'])
BGA_max_x = math.floor(BGA_layout['x_cp (mm)'].max() * 2) / 2
BGA_min_x = math.floor(BGA_layout['x_cp (mm)'].min() * 2) / 2
BGA_max_y = math.floor(BGA_layout['y_cp (mm)'].max() * 2) / 2
BGA_min_y = math.floor(BGA_layout['y_cp (mm)'].min() * 2) / 2

BGA_layout['Current_copy'] = pd.to_numeric(BGA_layout['Current / A']).abs()
def multiply_with_params_BGA(row, params_df):
    net = row['Net']
    current_copy = row['Current_copy']
    if net in params_df.columns:
        return current_copy * params_df.loc[0, net]
    return current_copy

# Multiply each row and store the results in a new column 'BGA_adjusted'
BGA_layout['BGA_adjusted'] = BGA_layout.apply(lambda row: multiply_with_params_BGA(row, df_parameters_input), axis=1)

# Find the minimum and maximum values in the 'BGA_adjusted' column
BGA_min_edc_adjusted = BGA_layout['BGA_adjusted'].min()
BGA_max_edc_adjusted = BGA_layout['BGA_adjusted'].max()

# Convert to float variables
BGA_min_edc_adjusted = float(BGA_min_edc_adjusted)
BGA_max_edc_adjusted = float(BGA_max_edc_adjusted)

# ratio calc.
# BGA_len_x = int(BGA_max_x - BGA_min_x)
# BGA_len_y = int(BGA_max_y - BGA_min_y)
# BGA_fraction = math.gcd(BGA_len_x, BGA_len_y)
# BGA_xratio = (BGA_len_x / BGA_fraction)
# BGA_yratio = (BGA_len_y / BGA_fraction)
# print('BGA fraction: ', BGA_fraction)
BGA_xratio = int(BGA_max_x - BGA_min_x)
BGA_yratio = int(BGA_max_y - BGA_min_y)


# ## define IMAX function ##

def process_file(current_file):
    power_rail_name = os.path.splitext(os.path.basename(current_file))[0].replace('_IMAX', '')

    print(power_rail_name)
    ## Data Cleaning & DataFrame Creation ##
    # desired text line
    target_line = "Vias"
    remove_keywords = "Voltage Probes"
    # boolean flag to track if the target line has been reached
    found = False
    # list to store extracted information
    data = []
    with open(current_file, "r") as f:
        for line in f:
            # set found to True if the target line is reached
            if target_line in line:
                break
        # add line to data list if the target line has been reached
        for line in f:
            include_line = True
            if remove_keywords in line:
                break
            if include_line and line.strip():
                data.append(line.strip().split('\t'))

    columns = data[0]
    data = data[1:]
    df = pd.DataFrame(data, columns=columns)
    # make all current value positive
    df["Current / A"] = df["Current / A"].astype(float).abs()

    # Bump Vias #
    col_title = [col for col in df_parameters_input.columns if power_rail_name in col]
    EDC_Current = df_parameters_input.loc[0, col_title].astype(float).iloc[0].round(3)
    Bump_Vias = df_parameters_input.loc[1, col_title].astype(str).iloc[0]

    # Extract rows with specific Bump Vias inputted by user
    df_bumpvias = df[df["Via"].str.contains(Bump_Vias)]

    # Select the columns "Net", "x (mm)", "y (mm)", and "Current / A" and store them in a new data frame
    df_bump_selected = df_bumpvias[["Via", "Net", "x (mm)", "y (mm)", "Current / A"]]

    # in nlargest, n = number of top ranking listed
    top_n_values = df_bump_selected['Current / A'].nlargest(3).tolist()

    # Use the current value stored in list 'top_n_values' as reference, create a new dataframe
    df_top_n = df_bump_selected[df_bump_selected['Current / A'].isin(top_n_values)].sort_values(by='Current / A',
                                                                                                ascending=False)
    df_top_n = df_top_n.rename(columns={'Via': 'Via Name', 'Current / A': 'Imax per 1A'}).drop('Net', axis=1)
    df_top_n['Imax(A) per EDC Current'] = df_top_n['Imax per 1A'].astype(float) * EDC_Current
    df_top_n
    Imax_bump = df_top_n['Imax(A) per EDC Current'].values[0].round(3)
    # Imax_bump

    df_bump_selected = df_bump_selected.rename(columns={'x (mm)': 'x', 'y (mm)': 'y', 'Current / A': 'Current'})
    df_bump_selected['x'] = round(df_bump_selected['x'].astype(float), 3)
    df_bump_selected['y'] = round(df_bump_selected['y'].astype(float), 3)
    df_bump_selected['Current'] = df_bump_selected['Current'].mul(EDC_Current)
    print(df_bump_selected['Current'])
    index = pd.MultiIndex.from_product(
        [range(int(bump_min_y), int(bump_max_y) + 1), range(int(bump_min_x), int(bump_max_x) + 1)],
        names=['y_coords', 'x_coords'])

    # pivot the dataframe
    x_coords = np.arange(start=bump_min_x, stop=bump_max_x + 0.5, step=0.5)
    y_coords = np.arange(start=bump_min_y, stop=bump_max_y + 0.5, step=0.5)
    # create an empty dataframe with the fixed x and y coordinates
    df_bump_plot = pd.DataFrame(
        {'x': x_coords.repeat(len(y_coords)), 'y': np.tile(y_coords, len(x_coords)), 'Current': np.nan})

    # fill in the corresponding Current values from df_bump_selected
    for index, row in df_bump_selected.iterrows():
        x_idx = int((row['x'] - min(x_coords)) / 0.5)
        y_idx = int((row['y'] - min(y_coords)) / 0.5)
        df_bump_plot.iloc[x_idx * len(y_coords) + y_idx, 2] = row['Current']

    # pivot the dataframe from long to wide form
    df_result_bump = df_bump_plot.pivot(index='y', columns='x', values='Current')
    df_result_bump.sort_index(level=0, ascending=False, inplace=True)

    # Convert wide-form DataFrame to long-form
    df_result_bump_long = df_result_bump.reset_index().melt(id_vars='y', var_name='x', value_name='Current')
    # Create dot heatmap using seaborn scatterplot
    fig, ax = plt.subplots(figsize=(bump_xratio, bump_yratio))
    ax.set_aspect('equal', adjustable='box')
    norm = mcolors.Normalize(vmin=min_edc_adjusted, vmax=max_edc_adjusted)
    # sns.scatterplot(data=df_result_bump_long, x='x', y='y', hue='Current', size='Current', sizes=(50, 50), palette='coolwarm', ax=ax, vmin=min_edc_adjusted, vmax=max_edc_adjusted)
    sns.scatterplot(data=df_result_bump_long, x='x', y='y', hue='Current', hue_norm=norm, sizes=(50, 50),
                    palette='coolwarm', ax=ax)

    sm = plt.cm.ScalarMappable(cmap='coolwarm', norm=norm)
    sm.set_array([])
    cbar = plt.colorbar(sm, ax=ax, shrink=0.3)

    # Configure the x-axis and y-axis labels
    ax.set_xticks(np.arange(bump_min_x, bump_max_x + 0.5, 0.5))
    ax.set_yticks(np.arange(bump_max_y, bump_min_y - 0.5, -0.5))
    ax.set_xticklabels([str(round(x, 1)) for x in np.arange(bump_min_x, bump_max_x + 0.5, 0.5)])
    ax.set_yticklabels([str(round(y, 1)) for y in np.arange(bump_max_y, bump_min_y - 0.5, -0.5)])
    leg = ax.legend(title='Current', bbox_to_anchor=(1.05, 1), loc='upper left', borderaxespad=0)
    plt.setp(ax.get_xticklabels(), rotation=90)
    plt.title(power_rail_name + " Bump Sum of Current Map")
    plt.savefig(power_rail_name + " bump_dot_heatmap.png", bbox_inches='tight')

    # display plot
    # plt.show()

    ## BGA ##
    BGA_Via = df_parameters_input.loc[2, col_title].astype(str).iloc[0]

    # Extract rows with specific BGA Vias inputted by user
    df_BGAvias = df[df["Via"].str.contains(BGA_Via)]
    # Select the columns "Net", "x (mm)", "y (mm)", and "Current / A" and store them in a new data frame
    df_BGA_selected = df_BGAvias[["Via", "Net", "x (mm)", "y (mm)", "Current / A"]]

    # in nlargest, n = number of top ranking listed
    BGA_top_n_values = df_BGA_selected['Current / A'].nlargest(3).tolist()

    # Use the current value stored in list 'BGA_top_n_values' as reference, create a new dataframe
    df_top_n_BGA = df_BGA_selected[df_BGA_selected['Current / A'].isin(BGA_top_n_values)].sort_values(by='Current / A',
                                                                                                      ascending=False)
    df_top_n_BGA = df_top_n_BGA.rename(columns={'Via': 'Via Name', 'Current / A': 'Imax per 1A'}).drop('Net', axis=1)
    df_top_n_BGA['Imax(A) per EDC Current'] = df_top_n_BGA['Imax per 1A'].astype(float) * EDC_Current
    # df_top_n_BGA
    Imax_BGA = df_top_n_BGA['Imax(A) per EDC Current'].values[0].round(3)

    df_BGA_selected = df_BGA_selected.rename(columns={'x (mm)': 'x', 'y (mm)': 'y', 'Current / A': 'Current'})
    df_BGA_selected['x'] = df_BGA_selected['x'].astype(float)
    df_BGA_selected['y'] = df_BGA_selected['y'].astype(float)
    df_BGA_selected['Current'] = df_BGA_selected['Current'].mul(EDC_Current)

    index_BGA = pd.MultiIndex.from_product(
        [range(int(BGA_min_y), int(BGA_max_y) + 1), range(int(BGA_min_x), int(BGA_max_x) + 1)],
        names=['y_coords_BGA', 'x_coords_BGA'])

    # pivot the dataframe
    x_coords_BGA = np.arange(start=BGA_min_x, stop=BGA_max_x + 0.5, step=0.5)
    # np.round(x_coords,2).tolist()
    y_coords_BGA = np.arange(start=BGA_min_y, stop=BGA_max_y + 0.5, step=0.5)
    # np.round(y_coords,2).tolist()
    # create an empty dataframe with the fixed x and y coordinates
    df_BGA_plot = pd.DataFrame(
        {'x': x_coords_BGA.repeat(len(y_coords_BGA)), 'y': np.tile(y_coords_BGA, len(x_coords_BGA)), 'Current': np.nan})

    # fill in the corresponding Current values from df_BGA_selected
    for index, row in df_BGA_selected.iterrows():
        x_idx_BGA = int((row['x'] - min(x_coords_BGA)) / 0.5)
        y_idx_BGA = int((row['y'] - min(y_coords_BGA)) / 0.5)
        df_BGA_plot.iloc[x_idx_BGA * len(y_coords_BGA) + y_idx_BGA, 2] = row['Current']

    # pivot the dataframe from long to wide form
    df_result_BGA = df_BGA_plot.pivot(index='y', columns='x', values='Current')
    df_result_BGA.sort_index(level=0, ascending=False, inplace=True)

    # Convert wide-form DataFrame to long-form
    df_result_BGA_long = df_result_BGA.reset_index().melt(id_vars='y', var_name='x', value_name='Current')
    fig, ax = plt.subplots(figsize=(BGA_xratio, BGA_yratio))
    ax.set_aspect('equal', adjustable='box')
    norm = mcolors.Normalize(vmin=BGA_min_edc_adjusted, vmax=BGA_max_edc_adjusted)
    # sns.scatterplot(data=df_result_bump_long, x='x', y='y', hue='Current', size='Current', sizes=(50, 50), palette='coolwarm', ax=ax, vmin=min_edc_adjusted, vmax=max_edc_adjusted)
    sns.scatterplot(data=df_result_BGA_long, x='x', y='y', hue='Current', hue_norm=norm, sizes=(50, 50),
                    palette='coolwarm', ax=ax)

    sm = plt.cm.ScalarMappable(cmap='coolwarm', norm=norm)
    sm.set_array([])
    cbar = plt.colorbar(sm, ax=ax, shrink=0.3)

    # Configure the x-axis and y-axis labels
    ax.set_xticks(np.arange(BGA_min_x, BGA_max_x + 0.5, 0.5))
    ax.set_yticks(np.arange(BGA_max_y, BGA_min_y - 0.5, -0.5))
    ax.set_xticklabels([str(round(x, 1)) for x in np.arange(BGA_min_x, BGA_max_x + 0.5, 0.5)])
    ax.set_yticklabels([str(round(y, 1)) for y in np.arange(BGA_max_y, BGA_min_y - 0.5, -0.5)])
    leg = ax.legend(title='Current', bbox_to_anchor=(1.05, 1), loc='upper left', borderaxespad=0)
    plt.setp(ax.get_xticklabels(), rotation=90)
    plt.title(power_rail_name + " BGA Sum of Current Map")

    plt.savefig(power_rail_name + " BGA_dot_heatmap.png", bbox_inches='tight')

    # display plot
    # plt.show()

    ## Create empty dataframe for Imax data storage ##
    Imax_list = [power_rail_name, EDC_Current, Imax_bump, Imax_BGA]
    Imax_df.loc[len(Imax_df)] = Imax_list

    ## Export 2 hetamap dataframes to excel sheet
    with pd.ExcelWriter(power_rail_name + 'DC_Imax_report.xlsx') as writer:
        df_result_bump.to_excel(writer, sheet_name=power_rail_name + '_bump')
        df_result_BGA.to_excel(writer, sheet_name=power_rail_name + '_BGA')


# In[ ]:

# ## loop through each IMAX file

# In[38]:

for file_name in list_of_files_IMAX:
    current_file = file_name
    process_file(current_file)
Imax_df

# In[28]:


# dfi.export(Imax_df.style.hide(axis='index'), 'IMAX_table.png')
dfi.export(Imax_df, 'IMAX_table.png')
# dfi.export(Imax_df.style, 'IMAX_table.png', table_conversion='matplotlib', index=False)


# In[ ]:

# ## input DCR files

# In[29]:


# directory = 'C:\\Users\\kaiwning\\Desktop\\test2'
# List files ending with 'DCR.dc' and also present as a column name in df_parameters_input
list_of_files_DCR = [f for f in os.listdir(directory) if
                     f.endswith('DCR.dc') and os.path.splitext(f)[0].rstrip('_DCR') in column_names]
# Add the directory path to the filenames
list_of_files_DCR = [os.path.join(directory, f) for f in list_of_files_DCR]
# print(list_of_files_DCR)


DCR_df = pd.DataFrame(columns=['Power Rail', 'Current (A)', 'DCR (mOhm)', 'IR Drop (mV)'])


# ## define DCR function

def process_DCR_file(current_file):
    power_rail_name = os.path.splitext(os.path.basename(current_file))[0].replace('_DCR', '')
    target_line = "Current Sources"
    found = False
    data = []
    with open(file_name, "r") as f:
        for line in f:
            if target_line in line:
                break
        for line in f:
            include_line = True
            if include_line and line.strip():
                data.append(line.strip().split('\t'))

    columns = data[0]
    data = data[1:]
    df = pd.DataFrame(data, columns=columns)
    col_title = [col for col in df_parameters_input.columns if power_rail_name in col + '_DCR']
    EDC_Current = df_parameters_input.loc[0, col_title].astype(float).iloc[0].round(3)
    DCR_mOhm = abs(float(df.loc[0, 'Voltage / V'])) * 1000
    IR_Drop = DCR_mOhm * EDC_Current
    # IR_Drop=round(IR_Drop,2)

    IR_Drop = round(IR_Drop, 2)
    DCR_mOhm = round(DCR_mOhm, 3)
    DCR_list = [power_rail_name, EDC_Current, DCR_mOhm, IR_Drop]
    DCR_df.loc[len(DCR_df)] = DCR_list


# ## loop through each DCR file ##

# In[32]:


for file_name in list_of_files_DCR:
    current_file = file_name
    process_DCR_file(current_file)
DCR_df

# In[33]:


# dfi.export(DCR_df.style.hide(axis='index'), 'DCR_table.png')
dfi.export(DCR_df, 'DCR_table.png')
# dfi.export(DCR_df.style, 'DCR_table.png', table_conversion='matplotlib', index=False)

# Create a PowerPoint presentation
prs = Presentation()
# change slide sizes to Widescreen
slide_size = (16, 9)
prs.slide_width, prs.slide_height = Inches(slide_size[0]), Inches(slide_size[1])


# def px_to_inches(path):
#     im = Image.open(path)
#     width = im.width / im.info['dpi'][0]
#     height = im.height / im.info['dpi'][1]
#     return (width, height)
def px_to_inches(path):
    im = Image.open(path)
    dpi = im.info.get('dpi', (96, 96))  # Use a default DPI value of (96, 96) if not available
    width = im.width / dpi[0]
    height = im.height / dpi[1]
    return (width, height)



slide = prs.slides.add_slide(prs.slide_layouts[5])
# Assuming there's only one textbox in the slide
textbox = slide.shapes[0]
# Add text to the textbox
text_frame = textbox.text_frame
text_frame.clear()  # Clear any existing text
title_text = text_frame.paragraphs[0].add_run()
title_text.text = "IMAX Table"
text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.LEFT

img = px_to_inches('IMAX_table.png')

# add bump heatmap to every slide:
# slide = prs.slides.add_slide(prs.slide_layouts[5])
if Inches(slide_size[0] - img[0]) > 8 and Inches(slide_size[1] - img[1]) > 4.5:
    width = Inches(img[0] * 2)
    height = Inches(img[1] * 2)
    left = Inches(slide_size[0] - img[0] * 2) / 2
    top = Inches(slide_size[1] - img[1] * 2) / 2
elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) > 0:
    width = Inches(14)
    height = Inches(img[1] * 14 / img[0])
    left = Inches(slide_size[0] - 14) / 2
    top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
elif Inches(slide_size[1] - img[1]) < 0 and Inches(slide_size[0] - img[0]) > 0:
    width = Inches(img[0] * 8 / img[1])
    height = Inches(8)
    left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
    top = Inches(slide_size[1] - 8) / 2
elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) < 0:
    if img[0] > img[1]:
        width = Inches(14)
        height = Inches(img[1] * 14 / img[0])
        left = Inches(slide_size[0] - 14) / 2
        top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
    else:
        width = Inches(img[0] * 8 / img[1])
        height = Inches(8)
        left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
        top = Inches(slide_size[1] - 8) / 2
else:
    width = Inches(img[0])
    height = Inches(img[1])
    left = Inches(slide_size[0] - img[0]) / 2
    top = Inches(slide_size[1] - img[1]) / 2
slide.shapes.add_picture('IMAX_table.png', left, top, width = width, height = height)

slide = prs.slides.add_slide(prs.slide_layouts[5])
textbox = slide.shapes[0]
# Add text to the textbox
text_frame = textbox.text_frame
text_frame.clear()  # Clear any existing text
title_text = text_frame.paragraphs[0].add_run()
title_text.text = "DCR Table"
text_frame.paragraphs[0].alignment = pptx.enum.text.PP_ALIGN.LEFT
img = px_to_inches('DCR_table.png')
# add bump heatmap to every slide:
# slide = prs.slides.add_slide(prs.slide_layouts[5])
if Inches(slide_size[0] - img[0]) > 8 and Inches(slide_size[1] - img[1]) > 4.5:
    width = Inches(img[0] * 2)
    height = Inches(img[1] * 2)
    left = Inches(slide_size[0] - img[0] * 2) / 2
    top = Inches(slide_size[1] - img[1] * 2) / 2
elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) > 0:
    width = Inches(14)
    height = Inches(img[1] * 14 / img[0])
    left = Inches(slide_size[0] - 14) / 2
    top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
elif Inches(slide_size[1] - img[1]) < 0 and Inches(slide_size[0] - img[0]) > 0:
    width = Inches(img[0] * 8 / img[1])
    height = Inches(8)
    left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
    top = Inches(slide_size[1] - 8) / 2
elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) < 0:
    if img[0] > img[1]:
        width = Inches(14)
        height = Inches(img[1] * 14 / img[0])
        left = Inches(slide_size[0] - 14) / 2
        top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
    else:
        width = Inches(img[0] * 8 / img[1])
        height = Inches(8)
        left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
        top = Inches(slide_size[1] - 8) / 2
else:
    width = Inches(img[0])
    height = Inches(img[1])
    left = Inches(slide_size[0] - img[0]) / 2
    top = Inches(slide_size[1] - img[1]) / 2
slide.shapes.add_picture('DCR_table.png', left, top, width = width, height = height)
print(slide_size)
print(img)

for f in list_of_files_IMAX:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    power_rail_name = os.path.splitext(os.path.basename(f))[0].replace('_IMAX', '')
    img = px_to_inches(power_rail_name + ' bump_dot_heatmap.png')
    print("Image dimensions in inches:", img)
    # add bump heatmap to every slide:
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    if slide_size[0] - img[0] > 8 and slide_size[1] - img[1] > 4.5:
        width = Inches(img[0] * 2)
        height = Inches(img[1] * 2)
        left = Inches(slide_size[0] - img[0] * 2) / 2
        top = Inches(slide_size[1] - img[1] * 2) / 2
    elif slide_size[0] - img[0] < 0 and slide_size[1] - img[1] > 0:
        width = Inches(14)
        height = Inches(img[1] * 14 / img[0])
        left = Inches(slide_size[0] - 14) / 2
        top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
    elif slide_size[1] - img[1] < 0 and slide_size[0] - img[0] > 0:
        width = Inches(img[0] * 8 / img[1])
        height = Inches(8)
        left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
        top = Inches(slide_size[1] - 8) / 2
    elif slide_size[0] - img[0] < 0 and slide_size[1] - img[1] < 0:
        if img[0] > img[1]:
            width = Inches(14)
            height = Inches(img[1] * 14 / img[0])
            left = Inches(slide_size[0] - 14) / 2
            top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
        else:
            width = Inches(img[0] * 8 / img[1])
            height = Inches(8)
            left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
            top = Inches(slide_size[1] - 8) / 2
    else:
        width = Inches(img[0])
        height = Inches(img[1])
        left = Inches(slide_size[0] - img[0]) / 2
        top = Inches(slide_size[1] - img[1]) / 2
    slide.shapes.add_picture(power_rail_name + ' bump_dot_heatmap.png', left, top, width=width, height=height)

    print("First condition:", (slide_size[0] - img[0]) > 8, (slide_size[1] - img[1]) > 4.5)
    print("Second condition:", (slide_size[0] - img[0]) < 0, (slide_size[1] - img[1]) > 0)
    print("Third condition:", (slide_size[1] - img[1]) < 0, (slide_size[0] - img[0]) > 0)
    print("Fourth condition:", (slide_size[0] - img[0]) < 0, (slide_size[1] - img[1]) < 0)

    ##add BGA heatmap to every slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    img = px_to_inches(power_rail_name + ' BGA_dot_heatmap.png')
    # add bump heatmap to every slide:
    # slide = prs.slides.add_slide(prs.slide_layouts[5])
    if Inches(slide_size[0] - img[0]) > 8 and Inches(slide_size[1] - img[1]) > 4.5:
        width = Inches(img[0] * 2)
        height = Inches(img[1] * 2)
        left = Inches(slide_size[0] - img[0] * 2) / 2
        top = Inches(slide_size[1] - img[1] * 2) / 2
    elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) > 0:
        width = Inches(14)
        height = Inches(img[1] * 14 / img[0])
        left = Inches(slide_size[0] - 14) / 2
        top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
    elif Inches(slide_size[1] - img[1]) < 0 and Inches(slide_size[0] - img[0]) > 0:
        width = Inches(img[0] * 8 / img[1])
        height = Inches(8)
        left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
        top = Inches(slide_size[1] - 8) / 2
    elif Inches(slide_size[0] - img[0]) < 0 and Inches(slide_size[1] - img[1]) < 0:
        if img[0] > img[1]:
            width = Inches(14)
            height = Inches(img[1] * 14 / img[0])
            left = Inches(slide_size[0] - 14) / 2
            top = Inches(slide_size[1] - img[1] * 14 / img[0]) / 2
        else:
            width = Inches(img[0] * 8 / img[1])
            height = Inches(8)
            left = Inches(slide_size[0] - img[0] * 8 / img[1]) / 2
            top = Inches(slide_size[1] - 8) / 2
    else:
        width = Inches(img[0])
        height = Inches(img[1])
        left = Inches(slide_size[0] - img[0]) / 2
        top = Inches(slide_size[1] - img[1]) / 2
    slide.shapes.add_picture(power_rail_name + ' BGA_dot_heatmap.png', left, top, width = width, height = height)

prs.save("DC_IMAX_report_test2.pptx")
os.startfile("DC_IMAX_report_test2.pptx")